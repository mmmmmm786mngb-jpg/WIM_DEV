#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Сборка презентации PPTX (16:9) для бизнес-заказчика по проекту выноса расчета СЧА/РСА.

Оформление повторяет HTML-версию: фирменная палитра ВИМ (взята из логотипа),
логотип на каждом слайде, те же 16 слайдов и тот же текст.

Координаты задаются в пикселях сетки 1280 x 720 (как в HTML) и переводятся
в дюймы: 1280 px = 13.333 in, то есть 96 px на дюйм. Размеры шрифта задаются
в пикселях и переводятся в пункты по формуле pt = px * 0.75.
"""

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

BASE = Path(__file__).resolve().parent.parent.parent / "Документация" / "02_scha_rsa_group3"
LOGO = BASE / "vim_logo.png"
OUT = BASE / "presentation_scha_rsa.pptx"

FONT = "Segoe UI"
MONO = "Consolas"

BRAND = RGBColor(0x52, 0x64, 0x8A)
BRAND_DARK = RGBColor(0x33, 0x40, 0x5C)
BRAND_DEEP = RGBColor(0x23, 0x2C, 0x41)
ACCENT = RGBColor(0x8D, 0xB8, 0xC7)
ACCENT_SOFT = RGBColor(0xDC, 0xE8, 0xEE)
TINT = RGBColor(0xA8, 0xB1, 0xC4)
GOLD = RGBColor(0xB9, 0x8A, 0x22)
GOLD_SOFT = RGBColor(0xFB, 0xF4, 0xE2)
GOLD_INK = RGBColor(0x6E, 0x53, 0x13)
GREEN = RGBColor(0x2E, 0x7D, 0x64)
GREEN_SOFT = RGBColor(0xE7, 0xF2, 0xEE)
GREEN_INK = RGBColor(0x24, 0x60, 0x4F)
RED = RGBColor(0xA8, 0x42, 0x3F)
RED_SOFT = RGBColor(0xF8, 0xEC, 0xEB)
INK = RGBColor(0x3A, 0x43, 0x58)
MUTED = RGBColor(0x6A, 0x74, 0x88)
LINE = RGBColor(0xDC, 0xE1, 0xE9)
BG = RGBColor(0xF3, 0xF5, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_TXT = RGBColor(0xC9, 0xD3, 0xE0)

TOTAL = 16


def safe_print(text):
    """Безопасный вывод в консоль Windows - только ASCII."""
    try:
        print(text)
    except UnicodeEncodeError:
        print(text.encode("ascii", "replace").decode("ascii"))


def px(v):
    """Перевод пикселей сетки 1280x720 в EMU."""
    return Emu(int(round(v / 96.0 * 914400)))


def fs(v):
    """Перевод размера шрифта из пикселей в пункты."""
    return Pt(v * 0.75)


# ---------------------------------------------------------------- примитивы


def rect(slide, x, y, w, h, fill=None, line=None, line_px=1, shape=MSO_SHAPE.RECTANGLE):
    """Прямоугольник без тени с заданной заливкой и рамкой."""
    sh = slide.shapes.add_shape(shape, px(x), px(y), px(w), px(h))
    sh.shadow.inherit = False
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = px(line_px)
    sh.text_frame.text = ""
    return sh


def box(slide, x, y, w, h):
    """Текстовый блок с обнуленными полями."""
    tb = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.paragraphs[0].text = ""
    return tf


def para(tf, markup, size=17, color=INK, bold=False, line=1.35, before=0, after=0,
         align=PP_ALIGN.LEFT, spacing=None, caps=False, first=None, font=FONT):
    """
    Добавляет абзац с поддержкой разметки: **жирный** и `код`.

    Параметр first=True использует уже существующий пустой первый абзац.
    """
    if first is None:
        first = not tf.paragraphs[0].runs and tf.paragraphs[0].text == "" and len(tf.paragraphs) == 1
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.line_spacing = line
    p.space_before = Pt(before)
    p.space_after = Pt(after)

    for chunk, is_bold, is_code in _parse(markup):
        r = p.add_run()
        r.text = chunk
        r.font.size = fs(size)
        r.font.bold = bold or is_bold
        r.font.name = MONO if is_code else font
        r.font.color.rgb = BRAND_DARK if (is_bold and color is INK) else color
        if spacing is not None:
            _set_spacing(r, spacing)
        if caps:
            _set_caps(r)
    return p


def _parse(markup):
    """Разбор строки с **жирным** и `кодом` на последовательность фрагментов."""
    out = []
    buf = ""
    i = 0
    bold = False
    code = False
    while i < len(markup):
        if markup.startswith("**", i):
            if buf:
                out.append((buf, bold, code))
                buf = ""
            bold = not bold
            i += 2
            continue
        if markup[i] == "`":
            if buf:
                out.append((buf, bold, code))
                buf = ""
            code = not code
            i += 1
            continue
        buf += markup[i]
        i += 1
    if buf:
        out.append((buf, bold, code))
    return out or [("", False, False)]


def _set_spacing(run, pts):
    """Межбуквенный интервал (в пунктах) через прямую правку XML."""
    run.font._rPr.set("spc", str(int(pts * 100)))


def _set_caps(run):
    run.font._rPr.set("cap", "all")


# ---------------------------------------------------------------- каркас слайда


def new_slide(prs, bg=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if bg is not None:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = bg
    return slide


def logo(slide, x=1224 - 148, y=40, h=30, dark_chip=False):
    """Логотип ВИМ. На темном фоне выводится на белой подложке."""
    w = h * (528.0 / 126.0)
    if dark_chip:
        rect(slide, x - 14, y - 9, w + 28, h + 18, fill=WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    slide.shapes.add_picture(str(LOGO), px(x + (148 - w) if not dark_chip else x), px(y), height=px(h))


def page_num(slide, n, color=TINT):
    tf = box(slide, 1024, 676, 200, 20)
    para(tf, "%d / %d" % (n, TOTAL), size=13, color=color, align=PP_ALIGN.RIGHT, spacing=1)


def header(slide, eyebrow, title, n, title_size=34):
    """Шапка контентного слайда: надзаголовок, заголовок, линейка."""
    logo(slide)
    tf = box(slide, 62, 38, 900, 20)
    para(tf, eyebrow, size=14, color=TINT, bold=True, spacing=1.6, caps=True)
    tf = box(slide, 62, 62, 1080, 90)
    para(tf, title, size=title_size, color=BRAND_DARK, bold=True, line=1.16)
    y = 62 + (44 if title_size >= 30 else 36) * (2 if len(title) > 62 else 1)
    rect(slide, 62, y + 12, 74, 4, fill=ACCENT)
    page_num(slide, n)
    return y + 42


def note(slide, x, y, w, h, title, text, kind="info", size=17):
    """Врезка с цветной полосой слева."""
    palette = {
        "info": (ACCENT, ACCENT_SOFT, BRAND_DARK),
        "gold": (GOLD, GOLD_SOFT, GOLD_INK),
        "green": (GREEN, GREEN_SOFT, GREEN_INK),
        "red": (RED, RED_SOFT, RED),
    }[kind]
    rect(slide, x, y, w, h, fill=palette[1])
    rect(slide, x, y, 5, h, fill=palette[0])
    tf = box(slide, x + 20, y + 15, w - 40, h - 26)
    if title:
        para(tf, "**%s**" % title, size=size - 1, color=palette[2], line=1.3, after=3)
        para(tf, text, size=size, color=INK, line=1.4)
    else:
        para(tf, text, size=size, color=INK, line=1.4)


def card(slide, x, y, w, h, head, text, top_color=BRAND, fill=BG):
    rect(slide, x, y, w, h, fill=fill, line=LINE)
    rect(slide, x, y, w, 4, fill=top_color)
    tf = box(slide, x + 22, y + 22, w - 44, h - 40)
    para(tf, "**%s**" % head, size=19, color=BRAND_DARK, line=1.25, after=8)
    para(tf, text, size=16, color=INK, line=1.4)


def stat(slide, x, y, w, h, value, label, color=BRAND, top=None):
    rect(slide, x, y, w, h, fill=BG, line=LINE)
    if top:
        rect(slide, x, y, w, 4, fill=top)
    tf = box(slide, x + 14, y + int(h * 0.28), w - 28, 60)
    para(tf, value, size=44, color=color, bold=True, align=PP_ALIGN.CENTER, line=1.0, after=8)
    para(tf, label, size=15, color=MUTED, align=PP_ALIGN.CENTER, line=1.3)


def bullets(slide, x, y, w, items, size=17, gap=16):
    """Список с круглыми маркерами цвета акцента."""
    cy = y
    for it in items:
        rect(slide, x + 3, cy + 7, 8, 8, fill=ACCENT, shape=MSO_SHAPE.OVAL)
        tf = box(slide, x + 25, cy, w - 25, 60)
        para(tf, it, size=size, color=INK, line=1.4)
        chars_per_line = max(10, int((w - 25) / (size * 0.55)))
        lines = max(1, int(len(_plain(it)) / chars_per_line) + 1)
        cy += lines * int(size * 1.45) + gap
    return cy


def _plain(markup):
    return "".join(c for c, _, _ in _parse(markup))


def table(slide, x, y, w, rows, col_ratio, header_row=True, size=15.5, row_h=None):
    """Таблица с темной шапкой, без рамок, с чередованием фона строк."""
    n_rows = len(rows)
    n_cols = len(rows[0])
    height = px(row_h * n_rows if row_h else 60 * n_rows)
    gf = slide.shapes.add_table(n_rows, n_cols, px(x), px(y), px(w), height)
    tbl = gf.table
    tbl.first_row = header_row
    tbl.horz_banding = False
    total = float(sum(col_ratio))
    for i, ratio in enumerate(col_ratio):
        tbl.columns[i].width = px(w * ratio / total)
    for r, row in enumerate(rows):
        if row_h:
            tbl.rows[r].height = px(row_h)
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.margin_left = px(14)
            cell.margin_right = px(12)
            cell.margin_top = px(9)
            cell.margin_bottom = px(9)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            if r == 0 and header_row:
                cell.fill.fore_color.rgb = BRAND_DARK
                cell.text_frame.word_wrap = True
                para(cell.text_frame, val, size=15, color=WHITE, bold=True, line=1.25, first=True)
            else:
                cell.fill.fore_color.rgb = WHITE if r % 2 else RGBColor(0xF7, 0xF9, 0xFC)
                cell.text_frame.word_wrap = True
                para(cell.text_frame, val, size=size, color=INK, line=1.3, first=True)
    return tbl


# ---------------------------------------------------------------- слайды


def slide_01(prs):
    s = new_slide(prs)
    rect(s, 0, 0, 1280, 720, fill=RGBColor(0xF2, 0xF5, 0xF9))
    rect(s, 0, 0, 14, 720, fill=BRAND)
    s.shapes.add_picture(str(LOGO), px(96), px(120), height=px(52))
    tf = box(s, 96, 240, 900, 140)
    para(tf, "Ускорение вечерних регламентных операций розничного ДУ",
         size=50, color=BRAND_DARK, bold=True, line=1.14)
    tf = box(s, 96, 372, 880, 40)
    para(tf, "Вынос расчета СЧА и РСА в отдельную группу 3 регламентных операций", size=26, color=BRAND, line=1.25)
    rect(s, 96, 438, 340, 42, fill=GOLD)
    tf = box(s, 112, 450, 312, 24)
    para(tf, "**Предложение на согласование**", size=15, color=WHITE, spacing=0.6)
    tf = box(s, 96, 560, 700, 100)
    para(tf, "**Задача:** IMDEV-8663, доработка 2 из 2", size=16, color=MUTED, line=1.7)
    para(tf, "**Кому:** Кукушкин В.Ю., Юрашева Ю.Ю.", size=16, color=MUTED, line=1.7)
    para(tf, "**Система:** база ДУ; новая группа 3. \"Расчет СЧА/РСА\"", size=16, color=MUTED, line=1.7)
    para(tf, "**Дата:** 15 августа 2026", size=16, color=MUTED, line=1.7)
    page_num(s, 1)


def slide_02(prs):
    s = new_slide(prs)
    header(s, "Контекст", "О чем идет речь", 2)
    y, h = 224, 130
    card(s, 62, y, 366, h, "Что оптимизируем",
         "Группу регламентных операций **\"2. Вечерние операции\"** по розничному доверительному управлению")
    card(s, 457, y, 366, h, "Какой объем",
         "Порядка **22 000 договоров** в одном прогоне, и клиентская база продолжает расти")
    card(s, 852, y, 366, h, "Что считается тяжело",
         "Два показателя: **СЧА** — стоимость чистых активов и **РСА** — размер собственных активов")
    note(s, 62, 452, 1156, 130, "Важное уточнение по терминологии",
         "Название \"вечерние операции\" — историческое. Фактически эта группа запускается **днем и "
         "обрабатывает предыдущий рабочий день**. Значит, длительный прогон нагружает базу в рабочие "
         "часы и отодвигает момент готовности данных за вчера.", kind="gold")
    tf = box(s, 62, 676, 500, 20)
    para(tf, "Предложение по доработке. База ДУ", size=13, color=TINT)


def slide_03(prs):
    s = new_slide(prs)
    header(s, "Проблема", "Сегодня показатели считаются отдельно по каждому договору", 3)
    y = 176
    w, gap = 196, 22
    labels = ["Договор 1", "Договор 2", "Договор 3", "Договор 22 000"]
    x = 62
    for i, lb in enumerate(labels):
        rect(s, x, y, w, 78, fill=BG, line=LINE)
        tf = box(s, x + 10, y + 18, w - 20, 46)
        para(tf, lb, size=14.5, color=INK, align=PP_ALIGN.CENTER, line=1.25, after=2)
        para(tf, "**5 тяжелых запросов**", size=15, color=BRAND_DARK, align=PP_ALIGN.CENTER, line=1.25)
        sep = "+" if i < 2 else ("..." if i == 2 else "")
        if sep:
            tf = box(s, x + w, y + 28, gap, 24)
            para(tf, sep, size=18, color=TINT, bold=True, align=PP_ALIGN.CENTER)
        x += w + gap
    tf = box(s, x, y + 28, 22, 24)
    para(tf, "=", size=18, color=TINT, bold=True, align=PP_ALIGN.CENTER)
    rect(s, x + 22, y, 324, 78, fill=RED_SOFT, line=RGBColor(0xE3, 0xC4, 0xC2))
    tf = box(s, x + 32, y + 14, 304, 50)
    para(tf, "**110 000**", size=22, color=RED, align=PP_ALIGN.CENTER, line=1.2, after=2)
    para(tf, "обращений к базе данных", size=14.5, color=INK, align=PP_ALIGN.CENTER, line=1.2)

    bullets(s, 62, 296, 620, [
        "Каждое обращение — **не справочный запрос, а расчет остатков на конец дня**: счета "
        "бухгалтерского учета плюс около 17 регистров управленческого учета.",
        "Одна и та же по своей сути работа выполняется **22 000 раз**. База многократно перечитывает "
        "одни и те же таблицы.",
        "Расчет идет **в общей очереди операций** по договору, вперемешку с остальным закрытием.",
    ])
    stat(s, 888, 296, 330, 150, "22 000", "договоров в прогоне,\nи число растет")
    note(s, 62, 500, 1156, 96, "Из чего состоят эти пять запросов на каждый договор",
         "Параметры учетной политики договора, курсы валют, остатки по счетам бухгалтерского учета, "
         "остатки по регистрам управленческого учета и — при соответствующей настройке — справедливая "
         "стоимость незавершенных сделок.")


def slide_04(prs):
    s = new_slide(prs)
    header(s, "Почему сейчас", "Время прогона растет пропорционально числу договоров", 4)
    y, h = 232, 150
    card(s, 62, y, 366, h, "Рост клиентской базы",
         "Каждый новый договор добавляет прогону свою порцию времени. Зависимость линейная, потолка у нее нет.")
    card(s, 457, y, 366, h, "Нагрузка в рабочие часы",
         "Прогон идет днем, одновременно с работой пользователей: тяжелые запросы конкурируют с ними "
         "за ресурсы сервера.")
    card(s, 852, y, 366, h, "Позднее закрытие \"вчера\"",
         "Чем дольше прогон, тем позже готовы показатели за предыдущий рабочий день и все, что от них зависит.")
    note(s, 62, 462, 1156, 118, "Вывод",
         "Дело не только в том, что прогон долгий. Текущая схема расчета **ограничивает дальнейший рост "
         "розничного ДУ**: при сохранении подхода время будет увеличиваться вместе с числом договоров.",
         kind="red")


def slide_05(prs):
    s = new_slide(prs, bg=BRAND_DARK)
    rect(s, 0, 0, 1280, 720, fill=BRAND_DARK)
    logo(s, x=1224 - 131, y=38, h=28, dark_chip=True)
    tf = box(s, 96, 268, 700, 24)
    para(tf, "Ключевая часть", size=15, color=ACCENT, bold=True, spacing=2, caps=True)
    tf = box(s, 96, 306, 940, 80)
    para(tf, "Что мы предлагаем", size=52, color=WHITE, bold=True, line=1.15)
    tf = box(s, 96, 404, 900, 80)
    para(tf, "Два решения, которые требуют вашего согласования. Первое меняет **момент** расчета, "
             "второе — **способ** расчета.", size=23, color=LIGHT_TXT, line=1.4)
    page_num(s, 5, color=RGBColor(0x7C, 0x89, 0xA3))


def _prop_body(slide, x, y, w, h, items, accent=BRAND, fill=WHITE, title=None):
    """
    Карточка предложения: заголовок и пары "метка - текст".

    Все абзацы выводятся в один текстовый блок, чтобы PowerPoint сам управлял
    переносами и блоки не накладывались друг на друга.
    """
    rect(slide, x, y, w, h, fill=fill, line=LINE)
    rect(slide, x, y, 6, h, fill=accent)
    tf = box(slide, x + 26, y + 22, w - 52, h - 40)
    label_color = accent if accent is GOLD else BRAND
    first = True
    if title:
        para(tf, "**%s**" % title, size=21,
             color=GOLD_INK if accent is GOLD else BRAND_DARK, line=1.2, after=8, first=True)
        first = False
    for label, text in items:
        para(tf, label, size=12.5, color=label_color, bold=True, spacing=0.8, caps=True,
             before=0 if first else 13, after=4, first=first)
        first = False
        para(tf, text, size=16, color=INK, line=1.42, first=False)


def slide_06(prs):
    s = new_slide(prs)
    header(s, "Предложение 1 из 2", "Вынести расчет показателей в отдельную группу 3", 6)
    y = 200
    _prop_body(s, 62, y, 700, 440, [
        ("Что меняем",
         "Сейчас этапов закрытия два: \"1. Утренние операции\" и \"2. Вечерние операции\". Добавляем "
         "третий — **\"3. Расчет СЧА/РСА\"** — и переносим в него расчет показателей вместе со связанными "
         "контрольными операциями. Этап стартует после того, как все вечерние операции по всем договорам "
         "завершены."),
        ("Зачем это нужно",
         "Показатели считаются по остаткам, которые формируются в ходе самого закрытия: переоценки, "
         "начисления, резервы. Посчитать их сразу по всем договорам можно только тогда, когда все эти "
         "операции уже завершены. **Без отдельного этапа пакетный расчет технически невозможен.**"),
    ])
    _prop_body(s, 788, y, 430, 440, [
        ("Появление третьего этапа",
         "В списке групп регламентных операций добавится этап \"3. Расчет СЧА/РСА\"."),
        ("Третий документ плана за день",
         "По договору будет три документа \"План регламентных операций ДУ\" вместо двух: утро, вечер, "
         "расчет показателей."),
        ("Порядок запуска",
         "Предлагаем автоматизировать: для оператора это остается одним действием."),
    ], accent=GOLD, fill=RGBColor(0xFE, 0xFC, 0xF6), title="Что именно согласовываем")


def slide_07(prs):
    s = new_slide(prs)
    header(s, "Предложение 2 из 2", "Считать показатели пакетно, по группам однотипных договоров", 7)
    y = 200
    _prop_body(s, 62, y, 780, 440, [
        ("Что меняем",
         "Вместо отдельного расчета по каждому договору система считает **сразу по группе договоров с "
         "одинаковыми настройками учетной политики**, а затем раскладывает готовый результат по договорам."),
        ("Почему группами, а не все одним запросом",
         "Состав данных зависит от параметров договора: ведется ли обособленный учет задолженности "
         "эмитента, как учитываются сделки РЕПО, к какому типу клиента относится портфель. Договоры с "
         "одинаковым набором параметров считать вместе можно, с разным — нельзя: результат исказится."),
        ("Что это дает",
         "Розничные портфели настроены однотипно, поэтому вместо 22 000 отдельных расчетов система "
         "выполнит **несколько групповых запросов**. Логика расчета остается той же — меняется только "
         "объем данных за одно обращение."),
    ])
    stat(s, 872, y, 346, 128, "110 000", "обращений к базе\nпо этому блоку сейчас", color=RED, top=RED)
    stat(s, 872, y + 146, 346, 128, "десятки", "обращений к базе\nпосле доработки", color=GREEN, top=GREEN)
    note(s, 872, y + 292, 346, 148, "Что согласовываем",
         "Подход и критерий приемки: суммы должны совпасть с текущими **полностью, до копейки**.",
         kind="green", size=16)


def _fbox(slide, x, y, w, h, head, text, kind="plain"):
    fill, line = {
        "plain": (BG, LINE),
        "heavy": (RED_SOFT, RGBColor(0xE3, 0xC4, 0xC2)),
        "good": (GREEN_SOFT, RGBColor(0xC3, 0xDE, 0xD4)),
    }[kind]
    rect(slide, x, y, w, h, fill=fill, line=line)
    tf = box(slide, x + 12, y + 16, w - 24, h - 26)
    para(tf, "**%s**" % head, size=15, color=BRAND_DARK, align=PP_ALIGN.CENTER, line=1.25, after=3)
    para(tf, text, size=14.5, color=INK, align=PP_ALIGN.CENTER, line=1.3)


def _chev(slide, x, y, w=26):
    tf = box(slide, x, y, w, 24)
    para(tf, ">", size=18, color=TINT, bold=True, align=PP_ALIGN.CENTER)


def slide_08(prs):
    s = new_slide(prs)
    header(s, "Схема", "Как есть и как будет", 8)
    tf = box(s, 62, 178, 900, 20)
    para(tf, "Как есть: расчет внутри общей очереди по каждому договору",
         size=15, color=MUTED, bold=True, spacing=1, caps=True)
    y = 210
    w = 176
    x = 62
    seq = [("Договор 1", "операции закрытия", "plain"), ("Договор 1", "расчет СЧА/РСА", "heavy"),
           ("Договор 2", "операции закрытия", "plain"), ("Договор 2", "расчет СЧА/РСА", "heavy")]
    for i, (h1, h2, kind) in enumerate(seq):
        _fbox(s, x, y, w, 78, h1, h2, kind)
        x += w
        if i < 3:
            _chev(s, x, y + 28)
            x += 26
    tf = box(s, x, y + 28, 40, 24)
    para(tf, "...", size=18, color=TINT, bold=True, align=PP_ALIGN.CENTER)
    _fbox(s, x + 40, y, 168, 78, "x 22 000", "договоров", "heavy")

    tf = box(s, 62, 336, 900, 20)
    para(tf, "Как будет: два этапа, расчет один раз на всех",
         size=15, color=MUTED, bold=True, spacing=1, caps=True)
    y = 368
    w2 = 356
    _fbox(s, 62, y, w2, 92, "Этап 2. Вечерние операции",
          "все операции закрытия по всем 22 000 договоров, как сейчас")
    _chev(s, 62 + w2, y + 34)
    _fbox(s, 62 + w2 + 26, y, w2, 92, "Этап 3. Расчет СЧА/РСА",
          "групповой расчет показателей сразу по всем договорам", "good")
    _chev(s, 62 + 2 * w2 + 26, y + 34)
    _fbox(s, 62 + 2 * w2 + 52, y, w2, 92, "Запись документов",
          "те же документы по каждому договору, с той же датой")

    note(s, 62, 500, 1156, 96, None,
         "Порядок операций внутри этапа сохраняется, зависимости между расчетом показателей и "
         "контрольными операциями не нарушаются. Меняется только момент и способ получения данных.")


def slide_09(prs):
    s = new_slide(prs)
    header(s, "Границы изменений", "Что при этом не меняется", 9)
    tiles = [
        ("Методика расчета", "Формулы, состав счетов и регистров, правила исключения и вычитания — тот же алгоритм"),
        ("Суммы показателей", "Результат совпадает до копейки. Это критерий приемки работ"),
        ("Документы", "Те же документы по каждому договору, с той же датой и временем"),
        ("Справка-расчет", "Тот же состав и та же детализация по показателям"),
        ("Отчетность в Банк России", "Источник данных и его содержимое не меняются"),
        ("Права и интерфейсы", "Роли, права доступа и формы отчетов не затрагиваются"),
    ]
    w, h = 376, 118
    for i, (tt, td) in enumerate(tiles):
        x = 62 + (i % 3) * (w + 14)
        y = 190 + (i // 3) * (h + 14)
        rect(s, x, y, w, h, fill=GREEN_SOFT, line=RGBColor(0xC9, 0xDF, 0xD8))
        tf = box(s, x + 18, y + 18, w - 36, h - 30)
        para(tf, "**%s**" % tt, size=16.5, color=GREEN_INK, line=1.25, after=5)
        para(tf, td, size=14.5, color=RGBColor(0x3E, 0x5A, 0x52), line=1.35)
    note(s, 62, 470, 1156, 112, "Это техническая оптимизация, а не изменение учета",
         "Меняется **когда** и **каким запросом** система получает данные. Что она считает и что "
         "записывает — остается прежним.", kind="green")


def slide_10(prs):
    s = new_slide(prs)
    header(s, "Эксплуатация", "Что изменится в ежедневной работе", 10)
    table(s, 62, 186, 1156, [
        ["Что", "Как будет"],
        ["Список этапов закрытия", "Появится третий этап — \"3. Расчет СЧА/РСА\""],
        ["Порядок запуска", "Сначала вечерние операции, затем расчет показателей. Предлагаем автоматизировать, "
                            "чтобы оператор по-прежнему выполнял одно действие"],
        ["Документы плана операций", "Три документа по договору за день вместо двух. Ограничение уникальности "
                                     "документов это штатно допускает"],
        ["Контроль выполнения", "Работает как раньше. Пока этап не выполнен, операции видны в мониторе как "
                                "невыполненные"],
        ["Повторный прогон по договору", "При перезапуске вечерних операций показатели нужно пересчитать — "
                                         "предусмотрим это программно, не полагаясь на память оператора"],
    ], col_ratio=[30, 70], row_h=48)
    note(s, 62, 552, 1156, 90, "Безопасный режим отказа",
         "Если третий этап не запустить, система **не выдаст неверные цифры**: показатели просто не будут "
         "рассчитаны, и это сразу видно в мониторе как невыполненная операция.", kind="gold", size=16)


def slide_11(prs):
    s = new_slide(prs)
    header(s, "Эффект", "Что получаем и как это проверим", 11)
    y, h = 196, 124
    card(s, 62, y, 366, h, "Быстрее прогон",
         "Число тяжелых обращений к базе по блоку СЧА/РСА сокращается на порядок — с десятков тысяч до десятков",
         top_color=GREEN)
    card(s, 457, y, 366, h, "Меньше нагрузки днем",
         "Прогон меньше конкурирует с работой пользователей за ресурсы сервера в рабочие часы", top_color=GREEN)
    card(s, 852, y, 366, h, "Предсказуемость",
         "Время расчета показателей начинает слабо зависеть от числа договоров, база может расти дальше",
         top_color=GREEN)
    note(s, 62, 352, 570, 214, "Замер до начала разработки",
         "На первом этапе мы измерим фактическую долю расчета СЧА/РСА в общем времени прогона и покажем "
         "ожидаемый выигрыш в цифрах. **Если доля окажется небольшой — сообщим об этом и предложим "
         "направить усилия на другой участок.**")
    note(s, 648, 352, 570, 214, "Критерий приемки",
         "Полное совпадение сумм СЧА и РСА по всем договорам за контрольный день, до копейки, при "
         "сокращении времени этапа. Сверка \"до и после\" выполняется на копии продуктивной базы.",
         kind="green")


def slide_12(prs):
    s = new_slide(prs)
    header(s, "Управление риском", "Риски и как мы их закрываем", 12)
    table(s, 62, 186, 1156, [
        ["Риск", "Как закрываем"],
        ["Расхождение сумм после доработки",
         "Внедряем в два шага: сначала переносим расчет в конец **без изменения алгоритма** и убеждаемся "
         "в полном совпадении, только потом включаем пакетный расчет"],
        ["Третий этап не запущен",
         "Автоматический последовательный запуск и отражение невыполненной операции в мониторе. Расчет "
         "можно доформировать отдельно, не повторяя все закрытие"],
        ["Повторный прогон без пересчета показателей",
         "Программная связка: при перезапуске вечерних операций по договору показатели помечаются к пересчету"],
        ["Выигрыш меньше ожидаемого",
         "Замер выполняется до разработки, решение \"делаем или не делаем\" принимается по факту"],
        ["Потребуется откат",
         "Откат — это возврат операций в вечерний этап, то есть **изменение настройки**, а не переустановка "
         "доработки"],
    ], col_ratio=[34, 66], row_h=68)


def slide_13(prs):
    s = new_slide(prs)
    header(s, "План", "Четыре этапа с контролем результата на каждом", 13)
    steps = [
        ("1", "Замер и анализ",
         "Измеряем фактическую долю расчета СЧА/РСА в общем времени прогона, анализируем настройки "
         "договоров. Итог — цифра ожидаемого выигрыша и решение о продолжении"),
        ("2", "Перенос в отдельный этап",
         "Создаем третий этап и переносим в него операции, алгоритм расчета не трогаем. Итог — суммы "
         "совпадают с текущими полностью"),
        ("3", "Включение пакетного расчета",
         "Включаем групповой расчет по однотипным договорам. Итог — сокращение времени этапа при тех же суммах"),
        ("4", "Регресс и сверка",
         "Прогон и сверка на копии продуктива, итоговый замер. Итог — протоколы сверки и замеров, "
         "готовность к переносу на продуктив"),
    ]
    y = 230
    w = 289
    for i, (n, th, tb_) in enumerate(steps):
        x = 62 + i * w
        if i < 3:
            rect(s, x + 46, y + 19, w - 64, 2, fill=LINE)
        rect(s, x, y, 40, 40, fill=BRAND, shape=MSO_SHAPE.OVAL)
        tf = box(s, x, y + 10, 40, 24)
        para(tf, n, size=17, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        tf = box(s, x, y + 56, w - 22, 46)
        para(tf, "**%s**" % th, size=17, color=BRAND_DARK, line=1.28)
        tf = box(s, x, y + 104, w - 22, 150)
        para(tf, tb_, size=14.5, color=INK, line=1.42)
    note(s, 62, 496, 1156, 84, None,
         "Оценку трудозатрат представим после согласования подхода: объем работ по третьему этапу "
         "зависит от результатов анализа настроек договоров на первом этапе.")


def slide_14(prs):
    s = new_slide(prs)
    header(s, "Решение за вами", "Что требуется от бизнес-заказчика", 14)
    asks = [
        ("1", "Согласовать появление третьего этапа закрытия",
         "И, как следствие, третьего документа плана регламентных операций по договору за день"),
        ("2", "Подтвердить требование к запуску",
         "Автоматический последовательный запуск двух этапов одним действием оператора — либо два "
         "отдельных запуска по регламенту эксплуатации"),
        ("3", "Согласовать окно на копии продуктивной базы",
         "Для замеров и сверки результатов \"до и после\""),
        ("4", "Подтвердить критерий приемки",
         "Полное совпадение сумм СЧА и РСА по всем договорам за контрольный день при сокращении времени этапа"),
    ]
    y, h = 190, 96
    for n, at, ad in asks:
        rect(s, 62, y, 1156, h, fill=GOLD_SOFT, line=RGBColor(0xE8, 0xD9, 0xAE))
        rect(s, 62, y, 6, h, fill=GOLD)
        rect(s, 84, y + 32, 32, 32, fill=GOLD, shape=MSO_SHAPE.OVAL)
        tf = box(s, 84, y + 40, 32, 20)
        para(tf, n, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        tf = box(s, 132, y + 24, 1060, 60)
        para(tf, "**%s**" % at, size=18, color=GOLD_INK, line=1.25, after=4)
        para(tf, ad, size=15.5, color=RGBColor(0x4B, 0x45, 0x35), line=1.35)
        y += h + 14


def slide_15(prs):
    s = new_slide(prs)
    header(s, "Техническое приложение 1 из 2", "Суть решения для отдела разработки", 15)
    y, h = 196, 330
    rect(s, 62, y, 570, h, fill=BG, line=LINE)
    rect(s, 62, y, 570, 4, fill=BRAND)
    tf = box(s, 84, y + 24, 526, 40)
    para(tf, "**Часть A. Отдельная группа операций (второй проход)**", size=19, color=BRAND_DARK, line=1.25)
    tf = box(s, 84, y + 70, 526, 240)
    para(tf, "Создается группа `3. \"Расчет СЧА/РСА\"`, в нее переносятся пять видов операций: расчет СЧА, "
             "расчет РСА, доходность портфеля, контроль уменьшения стоимости, контроль уровня рисков.",
         size=16, color=INK, line=1.42, after=10)
    para(tf, "**Ключевой факт:** привязка вида операции к группе хранится в иерархии справочника "
             "`ВидыОперацийЗакрытияПериода`, а не в разрезе договоров. Перенос — это смена реквизита "
             "`Родитель` у пяти предопределенных элементов. Миграция настроек по 22 000 договоров не "
             "требуется, индивидуальная периодичность сохраняется.", size=16, color=INK, line=1.42)

    rect(s, 648, y, 570, h, fill=BG, line=LINE)
    rect(s, 648, y, 570, 4, fill=ACCENT)
    tf = box(s, 670, y + 24, 526, 40)
    para(tf, "**Часть D1. Группировка по сигнатуре параметров**", size=19, color=BRAND_DARK, line=1.25)
    tf = box(s, 670, y + 70, 526, 240)
    para(tf, "Тело запроса не переписывается. Отбор `ДоговорДУ = &ДоговорДУ` заменяется на "
             "`ДоговорДУ В (&ДоговорыГруппы)`, договор добавляется в выборку и группировку.",
         size=16, color=INK, line=1.42, after=10)
    para(tf, "Группы формируются по 9 параметрам учетной политики, влияющим на состав счетов и регистров. "
             "Теоретический максимум: 12 групп по БУ, 32 по УУ. Для розничных портфелей ожидается 1-3 "
             "группы. Расчет выполняется пачками договоров — ограничение по памяти.",
         size=16, color=INK, line=1.42)

    note(s, 62, 552, 1156, 96, None,
         "Готовые таблицы доезжают до операции штатным каналом: предподготовка данных перед циклом плюс "
         "седьмой параметр исполняемой процедуры. Многопоточность, обработка отказов и регистрация "
         "исполнения операций остаются без изменений.")


def slide_16(prs):
    s = new_slide(prs)
    header(s, "Техническое приложение 2 из 2", "Результаты проверки технической возможности", 16)
    table(s, 62, 186, 1156, [
        ["Проверка по исходным текстам конфигурации", "Результат"],
        ["Несколько планов операций на договор в день",
         "Допустимо штатно: уникальность контролируется в разрезе даты, договора **и группы операций**"],
        ["Сдвиг даты документов расчета",
         "Не происходит: дата привязана к окончанию рабочего дня, от момента расчета не зависит"],
        ["Влияние операций 1000-1030 на входные данные расчета",
         "Отсутствует: регистр бухгалтерии и регистры накопления УУ они не затрагивают, поэтому "
         "предподготовка второго прохода корректна"],
        ["Зависимость контрольных операций от результата расчета",
         "Есть, поэтому переносятся вместе. Порядок внутри группы сохраняется"],
        ["Отображение новой группы в интерфейсе, блокировки, план-факт",
         "Доработка не требуется, конфликтов нет"],
    ], col_ratio=[46, 54], row_h=50)
    note(s, 62, 540, 1156, 96, "Вывод",
         "Решение реализуемо **без изменения архитектуры механизма регламентных операций**. Правки "
         "локализованы в предподготовке данных и в расчетном методе менеджера документа. Выполняется "
         "расширением конфигурации.", kind="green", size=16)
    tf = box(s, 62, 676, 700, 20)
    para(tf, "IMDEV-8663, доработка 2. Подробное описание — proposal_scha_rsa_batch_calculation.html",
         size=13, color=TINT)


def build():
    if not LOGO.exists():
        safe_print("ERROR: logo not found: %s" % LOGO)
        return 1

    prs = Presentation()
    prs.slide_width = px(1280)
    prs.slide_height = px(720)

    for fn in (slide_01, slide_02, slide_03, slide_04, slide_05, slide_06, slide_07, slide_08,
               slide_09, slide_10, slide_11, slide_12, slide_13, slide_14, slide_15, slide_16):
        fn(prs)

    prs.save(str(OUT))
    safe_print("PPTX created: %s" % OUT.name)
    safe_print("slides: %d" % len(prs.slides._sldIdLst))
    safe_print("size on disk: %.1f KB" % (OUT.stat().st_size / 1024))
    return 0


if __name__ == "__main__":
    sys.exit(build())

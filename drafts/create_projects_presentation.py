#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Генератор презентации о проектах WIM_DEV
Собирает информацию о проектах из bases/ и создает PowerPoint презентацию
"""

import os
import glob
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Конфигурация
BASES_PATH = "c:/1c/Cursor_1c/WIM_DEV/bases"
OUTPUT_PATH = "c:/1c/Cursor_1c/WIM_DEV/drafts/WIM_DEV_Projects_Portfolio.pptx"

# Цветовая схема (корпоративные цвета)
COLOR_PRIMARY = RGBColor(0x1A, 0x23, 0x7E)      # Темно-синий
COLOR_SECONDARY = RGBColor(0x00, 0x96, 0xC7)    # Голубой
COLOR_ACCENT = RGBColor(0xFF, 0x6B, 0x35)       # Оранжевый
COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)         # Темно-серый
COLOR_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)        # Светло-серый


def get_project_info(base_name, project_name):
    """Собирает информацию о проекте"""
    project_path = os.path.join(BASES_PATH, base_name, "projects", project_name)
    
    if not os.path.exists(project_path):
        return None
    
    info = {
        "name": project_name,
        "base": base_name,
        "extensions": 0,
        "documentation": 0,
        "scripts": 0,
        "tests": 0,
        "has_readme": False,
        "files": []
    }
    
    # Подсчет файлов по категориям
    for root, dirs, files in os.walk(project_path):
        # Пропускаем скрытые папки
        dirs[:] = [d for d in dirs if not d.startswith('.')]
        
        for file in files:
            file_lower = file.lower()
            info["files"].append(file)
            
            # Расширения конфигурации
            if file_lower.endswith(('.xml', '.bsl', '.epf', '.erf')):
                info["extensions"] += 1
            # Документация
            elif file_lower.endswith(('.md', '.html', '.txt', '.docx')):
                info["documentation"] += 1
            # Скрипты
            elif file_lower.endswith(('.py', '.ps1', '.bat', '.sh')):
                info["scripts"] += 1
            # Тесты
            elif 'test' in file_lower or file_lower.endswith(('.json', '.yaml', '.yml')):
                info["tests"] += 1
    
    # Проверка наличия README
    readme_path = os.path.join(project_path, "README.md")
    info["has_readme"] = os.path.exists(readme_path)
    
    return info


def collect_all_projects():
    """Собирает информацию о всех проектах"""
    projects_data = {
        "bases": {},
        "total_projects": 0,
        "total_files": 0,
        "bases_list": []
    }
    
    # Получаем список баз
    bases = []
    for item in os.listdir(BASES_PATH):
        item_path = os.path.join(BASES_PATH, item)
        if os.path.isdir(item_path) and not item.startswith('.') and item != 'shared':
            bases.append(item)
    
    projects_data["bases_list"] = bases
    
    # Собираем проекты по каждой базе
    for base in bases:
        projects_path = os.path.join(BASES_PATH, base, "projects")
        if not os.path.exists(projects_path):
            continue
        
        base_projects = []
        for project in os.listdir(projects_path):
            project_path = os.path.join(projects_path, project)
            if os.path.isdir(project_path) and not project.startswith('.'):
                project_info = get_project_info(base, project)
                if project_info:
                    base_projects.append(project_info)
                    projects_data["total_projects"] += 1
                    projects_data["total_files"] += len(project_info["files"])
        
        projects_data["bases"][base] = base_projects
    
    return projects_data


def add_title_slide(prs, data):
    """Слайд 1: Титульный"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Фон
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLOR_PRIMARY
    background.line.fill.background()
    
    # Заголовок
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2), Inches(9), Inches(1.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "WIM_DEV"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER
    
    # Подзаголовок
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(3.5), Inches(9), Inches(1)
    )
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Portfolio of 1C Projects"
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    p.alignment = PP_ALIGN.CENTER
    
    # Дата
    date_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(6.5), Inches(9), Inches(0.5)
    )
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = datetime.now().strftime("%B %Y")
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    p.alignment = PP_ALIGN.CENTER


def add_overview_slide(prs, data):
    """Слайд 2: Общий обзор"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Заголовок слайда
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Infrastructure Overview"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    
    # Карточки с показателями
    metrics = [
        ("Bases", str(len(data["bases_list"])), COLOR_PRIMARY),
        ("Projects", str(data["total_projects"]), COLOR_SECONDARY),
        ("Files", str(data["total_files"]), COLOR_ACCENT),
    ]
    
    x_pos = 0.5
    for label, value, color in metrics:
        # Карточка
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos), Inches(1.5), Inches(2.8), Inches(2)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = color
        card.line.fill.background()
        
        # Значение
        value_box = slide.shapes.add_textbox(
            Inches(x_pos), Inches(1.8), Inches(2.8), Inches(1)
        )
        tf = value_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER
        
        # Подпись
        label_box = slide.shapes.add_textbox(
            Inches(x_pos), Inches(2.8), Inches(2.8), Inches(0.5)
        )
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER
        
        x_pos += 3.1
    
    # Список баз
    bases_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(4), Inches(9), Inches(2.5)
    )
    tf = bases_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Configuration Bases:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT
    
    bases_text = "\n".join([f"  • {base}" for base in data["bases_list"]])
    p = tf.add_paragraph()
    p.text = bases_text
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_TEXT
    p.space_before = Pt(12)


def add_base_slide(prs, base_name, projects):
    """Слайд для конкретной базы с проектами"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Заголовок
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Base: {base_name}"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    
    if not projects:
        # Нет проектов
        msg_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2), Inches(9), Inches(1)
        )
        tf = msg_box.text_frame
        p = tf.paragraphs[0]
        p.text = "No active projects (only template 'example')"
        p.font.size = Pt(16)
        p.font.italic = True
        p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
        return
    
    # Фильтруем реальные проекты (не example)
    real_projects = [p for p in projects if p["name"] != "example"]
    
    if not real_projects:
        msg_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2), Inches(9), Inches(1)
        )
        tf = msg_box.text_frame
        p = tf.paragraphs[0]
        p.text = "No active projects (only template 'example')"
        p.font.size = Pt(16)
        p.font.italic = True
        p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
        return
    
    # Создаем таблицу проектов
    rows = len(real_projects) + 1
    cols = 5
    table = slide.shapes.add_table(
        rows, cols, Inches(0.5), Inches(1.2), Inches(9), Inches(0.6 * rows)
    ).table
    
    # Заголовки таблицы
    headers = ["Project", "Extensions", "Docs", "Scripts", "Tests"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_PRIMARY
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        paragraph.font.bold = True
        paragraph.font.size = Pt(11)
    
    # Заполняем данные
    for idx, project in enumerate(real_projects, 1):
        table.cell(idx, 0).text = project["name"]
        table.cell(idx, 1).text = str(project["extensions"])
        table.cell(idx, 2).text = str(project["documentation"])
        table.cell(idx, 3).text = str(project["scripts"])
        table.cell(idx, 4).text = str(project["tests"])
        
        # Форматирование ячеек
        for col in range(cols):
            cell = table.cell(idx, col)
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(10)
            if col == 0:
                paragraph.font.bold = True


def add_architecture_slide(prs):
    """Слайд о структуре проектов"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Заголовок
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Project Structure"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    
    # Структура папок
    structure_items = [
        ("Extensions/", "Configuration extensions, modifications", COLOR_SECONDARY),
        ("Documentation/", "Specifications, requirements, work logs", RGBColor(0x9B, 0x59, 0xB6)),
        ("Testing/", "Test scenarios, reports, Python COM tests", COLOR_ACCENT),
        ("Scripts/", "Export/import, migrations, integrations", RGBColor(0x27, 0xAE, 0x60)),
    ]
    
    y_pos = 1.5
    for folder, desc, color in structure_items:
        # Папка
        folder_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y_pos), Inches(2.5), Inches(0.6)
        )
        folder_box.fill.solid()
        folder_box.fill.fore_color.rgb = color
        folder_box.line.fill.background()
        
        folder_text = slide.shapes.add_textbox(
            Inches(0.5), Inches(y_pos), Inches(2.5), Inches(0.6)
        )
        tf = folder_text.text_frame
        tf.paragraphs[0].text = folder
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Описание
        desc_box = slide.shapes.add_textbox(
            Inches(3.2), Inches(y_pos), Inches(6), Inches(0.6)
        )
        tf = desc_box.text_frame
        tf.paragraphs[0].text = desc
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.color.rgb = COLOR_TEXT
        
        y_pos += 0.9
    
    # Примечание
    note_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(5.5), Inches(9), Inches(1)
    )
    tf = note_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Key Principles:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT
    
    notes = [
        "• Source paths stored in source-path.txt (not copied to repo)",
        "• Base configurations remain unchanged; all work in Extensions",
        "• Python + COM for automated testing and integration"
    ]
    for note in notes:
        p = tf.add_paragraph()
        p.text = note
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_TEXT
        p.space_before = Pt(4)


def add_summary_slide(prs, data):
    """Заключительный слайд"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Фон
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLOR_PRIMARY
    background.line.fill.background()
    
    # Заголовок
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2), Inches(9), Inches(1)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Summary"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER
    
    # Статистика
    stats = [
        f"{len(data['bases_list'])} Configuration Bases",
        f"{data['total_projects']} Active Projects",
        f"{data['total_files']} Source Files",
    ]
    
    y_pos = 3.5
    for stat in stats:
        stat_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(y_pos), Inches(9), Inches(0.6)
        )
        tf = stat_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"  {stat}"
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        p.alignment = PP_ALIGN.CENTER
        y_pos += 0.7


def create_presentation():
    """Главная функция создания презентации"""
    print("Collecting project data...")
    data = collect_all_projects()
    
    print(f"Found {len(data['bases_list'])} bases with {data['total_projects']} projects")
    
    # Создаем презентацию
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Добавляем слайды
    print("Creating slides...")
    add_title_slide(prs, data)
    add_overview_slide(prs, data)
    
    # Слайды по базам
    for base in data["bases_list"]:
        projects = data["bases"].get(base, [])
        add_base_slide(prs, base, projects)
    
    add_architecture_slide(prs)
    add_summary_slide(prs, data)
    
    # Сохраняем
    print(f"Saving presentation to: {OUTPUT_PATH}")
    prs.save(OUTPUT_PATH)
    print("Done!")
    
    return OUTPUT_PATH


if __name__ == "__main__":
    create_presentation()
